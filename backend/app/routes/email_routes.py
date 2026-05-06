import asyncio
import os
import logging
import io
import base64
import re
import json
import time
from typing import Optional, Union
from fastapi import APIRouter, HTTPException, Request, Response
from pydantic import BaseModel, validator
from app.services.task_extractor import extract_tasks
from app.services.task_store import save_tasks, get_tasks, mark_task_complete, delete_task, get_setting, set_setting, clear_all_tasks, save_feedback, get_feedback_examples, save_task_feedback, get_task_feedback_examples, save_category_override, get_category_overrides
from app.services.ai_service import summarize_email, generate_reply as ai_generate_reply, _normalize_category
from app.services.gmail_service import (
    get_emails, AuthRequiredError, get_auth_url, exchange_code,
    get_gmail_service, get_user_info, send_email_direct, creds_to_cookie
)

router = APIRouter()
logger = logging.getLogger("mailpilot.routes")

class SettingsRequest(BaseModel):
    fetch_limit: str = "10"
    default_action: str = "draft"
    fetch_priority: str = "all"
    ai_detail_level: str = "medium"
    ai_tone: str = "professional"
    ai_persona: str = ""
    ai_details: str = ""
    university_mode: bool = False
    university_roll_number: str = ""
    university_semester: str = ""
    university_section: str = ""
    university_course: str = ""
    university_specialization: str = ""
    university_campus: str = ""
    university_ignore_other_sections: bool = False
    university_filter_mode: str = "balanced"
    accent_color: str = "#8b5cf6"

    @validator('fetch_limit')
    def validate_limit(cls, v):
        # Handle both string and numeric inputs from frontend
        val = str(v)
        return val if val.isdigit() and 1 <= int(val) <= 200 else "10"
        
    @validator('ai_detail_level')
    def validate_detail(cls, v):
        val = str(v or "medium").lower()
        return val if val in ["short", "medium", "detailed"] else "medium"

    @validator('ai_tone')
    def validate_tone(cls, v):
        val = str(v or "professional").lower()
        return val if val in ["professional", "casual", "urgent"] else "professional"

    @validator('university_filter_mode')
    def validate_university_filter_mode(cls, v):
        val = str(v or "balanced").lower()
        return val if val in ["balanced", "strict"] else "balanced"

    @validator(
        'ai_persona',
        'ai_details',
        'university_roll_number',
        'university_semester',
        'university_section',
        'university_course',
        'university_specialization',
        'university_campus'
    )
    def trim_long_text(cls, v):
        return str(v or "").strip()[:1600]

class AuthCallbackRequest(BaseModel):
    code: str
    state: Optional[str] = None
    code_verifier: Optional[str] = None

class AnalyzeRequest(BaseModel):
    id: str
    subject: str
    sender: str
    snippet: str
    retry: bool = False
    ai_detail_level: Optional[str] = None
    ai_persona: Optional[str] = None

class ReplyRequest(BaseModel):
    content: str
    tone: Optional[str] = None
    ai_details: Optional[str] = None

class ChatRequest(BaseModel):
    message: str

class BriefRequest(BaseModel):
    emails: list = []
    tasks: list = []

class DeepDiveRequest(BaseModel):
    id: str
    threadId: Optional[str] = None
    subject: str = ""
    sender: str = ""
    snippet: str = ""

@router.post("/chat-with-james")
async def chat_with_james_route(req: ChatRequest, request: Request):
    user_email = _get_user_email(request)
    from app.services.ai_service import chat_with_james
    from app.services.gmail_service import get_emails
    from app.services.task_store import get_tasks
    
    message = (req.message or "").strip()
    if not message:
        return {"reply": "Boss, send me a command and I'll jump on it."}

    lowered = message.lower()
    needs_email_context = any(word in lowered for word in [
        "email", "mail", "inbox", "unread", "recent", "thread", "message", "sender", "from"
    ])

    async def safe_get_tasks():
        if not user_email:
            return []
        try:
            return await asyncio.wait_for(asyncio.to_thread(get_tasks, user_email), timeout=2.5)
        except Exception as e:
            logger.warning(f"James task context unavailable: {e}")
            return []

    async def safe_get_emails():
        if not needs_email_context:
            return []
        try:
            return await asyncio.wait_for(asyncio.to_thread(get_emails, _get_cookie(request)), timeout=5.0)
        except Exception as e:
            logger.warning(f"James email context unavailable: {e}")
            return []
    
    emails_raw, tasks = await asyncio.gather(safe_get_emails(), safe_get_tasks())
    emails = emails_raw[:3]
    active_tasks = [t for t in tasks if not t.get('completed')][:8]
    
    context = {
        "recent_emails": [
            {
                "subject": e.get('subject', '(No Subject)')[:120],
                "sender": e.get('sender', 'Unknown')[:120],
                "snippet": e.get('snippet', '')[:220]
            }
            for e in emails
        ],
        "active_tasks": active_tasks,
        "email_context_loaded": needs_email_context
    }
    
    reply = await chat_with_james(message[:1000], context)
    return {"reply": reply}

@router.post("/morning-brief")
async def morning_brief_route(req: BriefRequest, request: Request):
    user_email = _get_user_email(request)
    server_tasks = get_tasks(user_email) if user_email else []
    tasks = (req.tasks or server_tasks)[:12]
    emails = (req.emails or [])[:12]
    from app.services.ai_service import generate_morning_brief
    async with AI_SEMAPHORE:
        brief = await asyncio.to_thread(generate_morning_brief, emails, tasks)
    return {"brief": brief}

@router.post("/deep-dive")
async def deep_dive_route(req: DeepDiveRequest, request: Request):
    attachments = []
    content = req.snippet or ""
    attachment_texts = []
    try:
        service = get_gmail_service(_get_cookie(request))
        message = service.users().messages().get(userId="me", id=req.id, format="full").execute()
        payload = message.get("payload", {})
        from app.services.gmail_service import extract_body
        content = extract_body(payload) or content

        def decode_attachment_data(data):
            if not data: return b""
            missing = len(data) % 4
            if missing:
                data += "=" * (4 - missing)
            return base64.urlsafe_b64decode(data.encode())

        def extract_attachment_text(filename, mime_type, raw_bytes):
            name = (filename or "").lower()
            try:
                if name.endswith((".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".xml", ".html", ".htm", ".ics", ".rtf")):
                    text = raw_bytes.decode("utf-8", errors="ignore")
                    if name.endswith(".html") or name.endswith(".htm") or mime_type == "text/html":
                        from bs4 import BeautifulSoup
                        text = BeautifulSoup(text, "html.parser").get_text("\n")
                    return text[:6000]
                if mime_type and mime_type.startswith("text/"):
                    return raw_bytes.decode("utf-8", errors="ignore")[:6000]
                if name.endswith(".pdf") or mime_type == "application/pdf":
                    from pypdf import PdfReader
                    reader = PdfReader(io.BytesIO(raw_bytes))
                    return "\n".join((page.extract_text() or "") for page in reader.pages[:5])[:6000]
                if name.endswith(".docx") or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    from docx import Document
                    doc = Document(io.BytesIO(raw_bytes))
                    return "\n".join(p.text for p in doc.paragraphs[:120])[:6000]
                if name.endswith(".pptx") or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(raw_bytes))
                    slide_text = []
                    for slide in prs.slides[:20]:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                    return "\n".join(slide_text)[:6000]
                if name.endswith(".xlsx") or mime_type in {
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "application/vnd.ms-excel"
                }:
                    from openpyxl import load_workbook
                    wb = load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
                    rows = []
                    for ws in wb.worksheets[:5]:
                        rows.append(f"[Sheet] {ws.title}")
                        for row in ws.iter_rows(max_row=40, max_col=12, values_only=True):
                            values = [str(cell).strip() for cell in row if cell not in (None, "")]
                            if values:
                                rows.append(" | ".join(values))
                    return "\n".join(rows)[:6000]
                if name.endswith(".odp"):
                    import zipfile
                    with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
                        content_xml = zf.read("content.xml").decode("utf-8", errors="ignore")
                    return re.sub(r"<[^>]+>", " ", content_xml)[:6000]
                if name.endswith(".ods"):
                    import zipfile
                    with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
                        content_xml = zf.read("content.xml").decode("utf-8", errors="ignore")
                    return re.sub(r"<[^>]+>", " ", content_xml)[:6000]
                if name.endswith(".odt"):
                    import zipfile
                    with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
                        content_xml = zf.read("content.xml").decode("utf-8", errors="ignore")
                    return re.sub(r"<[^>]+>", " ", content_xml)[:6000]
                if mime_type and mime_type.startswith("image/"):
                    from app.services.ai_service import ocr_image_with_nvidia
                    ocr_text = ocr_image_with_nvidia(raw_bytes, mime_type)
                    return ocr_text or "[Image attachment detected, but OCR model was unavailable or could not read text.]"
            except Exception as e:
                logger.warning(f"Attachment text extraction failed for {filename}: {e}")
            return ""

        def collect_parts(part):
            filename = part.get("filename")
            body = part.get("body", {})
            if filename:
                item = {
                    "filename": filename,
                    "mimeType": part.get("mimeType"),
                    "size": body.get("size", 0),
                    "attachmentId": body.get("attachmentId")
                }
                attachments.append(item)
                if body.get("attachmentId") and len(attachment_texts) < 3 and int(body.get("size", 0) or 0) <= 2_500_000:
                    try:
                        att = service.users().messages().attachments().get(
                            userId="me", messageId=req.id, id=body["attachmentId"]
                        ).execute()
                        raw = decode_attachment_data(att.get("data", ""))
                        extracted = extract_attachment_text(filename, part.get("mimeType"), raw)
                        if extracted:
                            attachment_texts.append(f"Attachment: {filename}\n{extracted}")
                    except Exception as e:
                        logger.warning(f"Attachment fetch failed for {filename}: {e}")
            for child in part.get("parts", []) or []:
                collect_parts(child)

        collect_parts(payload)
    except AuthRequiredError:
        return {"needs_auth": True}
    except Exception as e:
        logger.warning(f"Deep dive Gmail enrichment unavailable: {e}")

    from app.services.ai_service import generate_deep_dive
    user_email = _get_user_email(request)
    persona = _build_effective_persona(user_email) if user_email else ""
    async with AI_SEMAPHORE:
        analysis = await asyncio.to_thread(
            generate_deep_dive,
            req.subject,
            req.sender,
            f"{content}\n\nATTACHMENT TEXT:\n" + "\n\n".join(attachment_texts),
            attachments,
            persona
        )
    return {"analysis": analysis or "Boss, I could not complete the deep dive this time.", "attachments": attachments, "attachment_text_found": bool(attachment_texts)}

class FeedbackRequest(BaseModel):
    id: str
    is_positive: bool
    snippet: str
    summary: str

class TaskFeedbackRequest(BaseModel):
    is_positive: bool
    task_text: Optional[str] = None
    item_id: str

class CategoryOverrideRequest(BaseModel):
    item_id: str
    category: Optional[str] = None

    @validator('item_id')
    def validate_item_id(cls, v):
        value = str(v or "").strip()
        if not value or len(value) > 255:
            raise ValueError("Invalid item_id")
        return value

    @validator('category')
    def validate_category(cls, v):
        if v in (None, "", "null"):
            return None
        normalized = _normalize_category(str(v), "")
        if not normalized or len(normalized) > 32:
            raise ValueError("Invalid category")
        return normalized

class SendEmailRequest(BaseModel):
    to: str
    subject: str
    body: str

def _build_effective_persona(user_email: str, explicit_persona: Optional[str] = None):
    lines = []
    base = explicit_persona if explicit_persona is not None else get_setting(user_email, "ai_persona", "")
    if base:
        lines.append(base)

    if str(get_setting(user_email, "university_mode", "false")).lower() == "true":
        lines.append("UNIVERSITY MODE ACTIVE: apply academic relevance filtering.")
        fields = [
            ("Roll number", "university_roll_number"),
            ("Semester", "university_semester"),
            ("Section", "university_section"),
            ("Course", "university_course"),
            ("Specialization", "university_specialization"),
            ("Campus", "university_campus"),
        ]
        for label, key in fields:
            value = get_setting(user_email, key, "")
            if value:
                lines.append(f"{label}: {value}")

        filter_mode = get_setting(user_email, "university_filter_mode", "balanced")
        strict_mode = filter_mode == "strict"
        lines.append(
            "University filter mode: STRICT MATCH. If an academic email is clearly for a different section, roll number, semester, course, specialization, or campus, classify it as FILTERED_NOISE unless it also explicitly applies to me."
            if strict_mode
            else "University filter mode: BALANCED CAMPUS. Use academic profile strongly, but keep nearby section/course notices as STRATEGIC_FYI when they may still affect me."
        )
        lines.append("Keep Assignment Update, Quiz Notice, and Course Alert as separate categories when they fit.")

    return "\n".join(lines)

# Concurrency control for AI processing
AI_SEMAPHORE = asyncio.Semaphore(10)

# ─── Cookie Auth Helpers ───

def _get_cookie(request: Request):
    """Get the auth token cookie value."""
    return request.cookies.get("mailpilot_token")

def _cookie_secure():
    """Use secure cookies in production while allowing local HTTP development."""
    redirect_uri = os.getenv("GMAIL_REDIRECT_URI", "")
    return os.getenv("VERCEL") == "1" or redirect_uri.startswith("https://")

def _set_token_cookie(response: Response, creds, user_email=None):
    """Set the auth token as a secure HTTP cookie with optional cached user_email."""
    response.set_cookie(
        key="mailpilot_token", value=creds_to_cookie(creds, user_email),
        httponly=True, secure=_cookie_secure(), samesite="lax",
        max_age=60*60*24*30, path="/"
    )

def _get_user_email(request: Request):
    """Returns the authenticated user's email or None. Fast path: check signed cookie first."""
    cookie = _get_cookie(request)
    if not cookie: return None
    try:
        if '.' not in cookie: return None
        payload, signature = cookie.split('.', 1)
        
        # Verify signature for safety
        from app.config.settings import settings
        import hmac, hashlib
        expected_sig = hmac.new(settings.SESSION_SECRET.encode(), payload.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(signature, expected_sig):
            logger.error("Auth Security: Signature mismatch.")
            return None

        # Root Cause Fix: Handle missing padding in Base64 payload
        missing_padding = len(payload) % 4
        if missing_padding:
            payload += '=' * (4 - missing_padding)
            
        data = json.loads(base64.urlsafe_b64decode(payload).decode())
        email = data.get('user_email')
        if email: return email.lower().strip()
        return None
    except Exception as e: 
        logger.error(f"Critical Auth Parsing Error: {e}")
        return None

# ─── Shared Background Status Tracker ───
ANALYSIS_STATUS = {} # In-memory status for Hobby (cleared on cold start, but good for active sessions)

async def run_analysis_background(req: AnalyzeRequest, user_email: str, cookie: str, detail_level: str, persona: str):
    """Background worker for deep intelligence extraction."""
    job_id = f"{user_email}:{req.id}"
    ANALYSIS_STATUS[job_id] = {"status": "processing"}
    
    try:
        async with AI_SEMAPHORE:
            context = f"Subject: {req.subject}\nSender: {req.sender}\n\nBody:\n{req.snippet}"[:4000]
            # 10-Shot Feedback Learning
            pos_examples, neg_examples = await asyncio.to_thread(get_feedback_examples, user_email, 5) if user_email else ([], [])
            task_pos, task_neg = await asyncio.to_thread(get_task_feedback_examples, user_email, 5) if user_email else ([], [])
            
            tasks_task = asyncio.to_thread(extract_tasks, context, task_pos, task_neg, req.retry, persona)
            summary_task = asyncio.to_thread(summarize_email, context, detail_level, req.retry, pos_examples, neg_examples, user_email, persona)
            
            tasks, summary_data = await asyncio.gather(tasks_task, summary_task)
            
            summary = summary_data.get("summary", "No summary.") if isinstance(summary_data, dict) else str(summary_data)
            category = summary_data.get("category", "STRATEGIC_FYI") if isinstance(summary_data, dict) else "STRATEGIC_FYI"
            james_note = summary_data.get("james_note") if isinstance(summary_data, dict) else None
            
            final_tasks = await asyncio.to_thread(save_tasks, tasks, user_email, message_id=req.id, source_snippet=context) if user_email else tasks

            ANALYSIS_STATUS[job_id] = {
                "status": "complete",
                "result": {
                    "id": req.id, 
                    "summary": summary, 
                    "category": category, 
                    "tasks": final_tasks, 
                    "james_note": james_note, 
                    "was_corrected": summary_data.get("was_corrected") if isinstance(summary_data, dict) else False,
                    "engine": summary_data.get("engine", "Nano 8B") if isinstance(summary_data, dict) else "Nano 8B"
                }
            }
    except Exception as e:
        logger.error(f"Background Analysis Failed for {req.id}: {e}")
        ANALYSIS_STATUS[job_id] = {"status": "error", "message": str(e)}

# ─── Routes ───

@router.get("/user-info")
async def user_info(request: Request, response: Response):
    try:
        from app.services.gmail_service import authenticate_gmail
        from googleapiclient.discovery import build
        cookie = _get_cookie(request)
        creds, was_refreshed = authenticate_gmail(cookie)
        service = build("gmail", "v1", credentials=creds)
        info = get_user_info(service)
        if info: 
            # Force cookie refresh if Google rotated tokens or an older cookie lacks user_email.
            if was_refreshed or not _get_user_email(request):
                _set_token_cookie(response, creds, info.get('email'))
            return info
        raise HTTPException(status_code=401, detail="User not authenticated")
    except AuthRequiredError:
        return {"needs_auth": True}

# NOTE: /reinforcement-status removed to comply with 'Strictly on User Device' requirement.
# All alignment telemetry is now computed and stored in the frontend localStorage.

@router.get("/emails-raw")
async def get_emails_raw(request: Request, response: Response):
    try:
        from app.services.gmail_service import get_emails
        cookie = _get_cookie(request)
        emails = get_emails(cookie)
        
        user_email = _get_user_email(request)
        if user_email:
            from app.services.task_store import get_tasks
            from app.models import SessionLocal, Task, CategoryOverride
            db = SessionLocal()
            try:
                # Optimized: Fetch all relevant data for this user in bulk
                all_tasks = db.query(Task).filter(Task.user_email == user_email, Task.archived == False).all()
                overrides = {o.item_id: o.category for o in db.query(CategoryOverride).filter(CategoryOverride.user_email == user_email).all()}

                for email in emails:
                    eid = email.get('id')
                    
                    # Apply Category Overrides
                    if eid in overrides:
                        email['category'] = overrides[eid]
                        email['category_override'] = overrides[eid]
                    
                    # Apply Associated Tasks
                    email_tasks = [t for t in all_tasks if t.message_id == eid]
                    if email_tasks:
                        email['tasks'] = [
                            {"id": t.id, "task": t.task_text, "deadline": t.deadline, "priority": t.priority, "completed": t.completed}
                            for t in email_tasks
                        ]
            finally:
                db.close()
        
        return {"data": emails}
    except AuthRequiredError:
        return {"needs_auth": True}
    except Exception as e:
        logger.error(f"Error in emails-raw: {e}")
        raise HTTPException(status_code=500, detail="Failed to fetch emails")

@router.get("/auth-url")
async def auth_url():
    """Generate a new OAuth URL. Returns auth_url, state, and code_verifier."""
    try:
        result = get_auth_url()
        return result
    except Exception as e:
        logger.error(f"Auth URL error: {e}")
        raise HTTPException(status_code=500, detail=f"Auth Setup Failed: {str(e)}")

from fastapi import BackgroundTasks, Response
from fastapi.responses import JSONResponse

@router.post("/analyze-single")
async def analyze_single(req: AnalyzeRequest, request: Request, response: Response, background_tasks: BackgroundTasks):
    cookie = _get_cookie(request)
    user_email = _get_user_email(request)
    
    # We need user_email for the job_id
    if not user_email:
        from app.services.gmail_service import authenticate_gmail, get_user_info
        from googleapiclient.discovery import build
        creds, was_refreshed = authenticate_gmail(cookie)
        info = get_user_info(build("gmail", "v1", credentials=creds))
        user_email = (info.get('email') or "").lower().strip()
        if was_refreshed:
            _set_token_cookie(response, creds, user_email)

    job_id = f"{user_email}:{req.id}"
    
    # 1. CHECK IN-MEMORY STATUS: Avoid double-triggering
    status = ANALYSIS_STATUS.get(job_id, {})
    if status.get("status") == "processing":
        return JSONResponse(status_code=202, content={"id": req.id, "status": "processing"})
    if status.get("status") == "complete" and not req.retry:
        return status["result"]

    # 2. TRIGGER BACKGROUND TASK: Truly decouple the request to avoid Vercel timeouts
    detail_level = req.ai_detail_level or get_setting(user_email, "ai_detail_level", "medium")
    persona = _build_effective_persona(user_email, req.ai_persona)
    
    background_tasks.add_task(run_analysis_background, req, user_email, cookie, detail_level, persona)
    
    return JSONResponse(
        status_code=202,
        content={"id": req.id, "status": "processing", "message": "James is analyzing in the background."}
    )


@router.get("/analyze/status/{message_id}")
async def get_analysis_status(message_id: str, request: Request):
    user_email = _get_user_email(request)
    if not user_email: raise HTTPException(status_code=401)
    
    job_id = f"{user_email}:{message_id}"
    status = ANALYSIS_STATUS.get(job_id, {"status": "not_started"})
    
    if status["status"] == "complete":
        return status["result"]
    return status

@router.post("/feedback")
def submit_feedback(req: FeedbackRequest, request: Request):
    user_email = _get_user_email(request)
    if not user_email:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    result = save_feedback(user_email, req.is_positive, req.snippet[:2000], req.summary, req.id)
    if not result:
        raise HTTPException(status_code=500, detail="Failed to save feedback")
    return {"status": "success", "action": result}

@router.post("/task-feedback")
def submit_task_feedback(req: TaskFeedbackRequest, request: Request):
    user_email = _get_user_email(request)
    if not user_email:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    # item_id is the task database ID
    try:
        tid = int(req.item_id)
        result = save_task_feedback(user_email, req.is_positive, tid)
    except (ValueError, TypeError):
        logger.error(f"Invalid item_id in task-feedback: {req.item_id}")
        raise HTTPException(status_code=400, detail="Invalid task ID")

    if not result:
        raise HTTPException(status_code=500, detail="Failed to save task feedback")
    return {"status": "success", "action": result}

@router.post("/category-override")
def category_override_route(req: CategoryOverrideRequest, request: Request):
    user_email = _get_user_email(request)
    if not user_email:
        raise HTTPException(status_code=401, detail="Unauthorized")
    result = save_category_override(user_email, req.item_id, req.category)
    if not result:
        raise HTTPException(status_code=500, detail="Failed to save category override")
    return {"status": "success", "action": result}

@router.post("/coo/consult-task")
async def consult_coo_on_task(req: dict, request: Request):
    """Escalates a task to the COO (340B) for a strategic operational roadmap."""
    task_id = req.get("task_id")
    user_email = _get_user_email(request)
    if not user_email:
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    from app.services.task_store import get_tasks
    from app.services.ai_service import generate_ai_response, COO_SYSTEM_PROMPT
    
    # Fetch task context
    all_tasks = get_tasks(user_email)
    target_task = next((t for t in all_tasks if str(t['id']) == str(task_id)), None)
    
    if not target_task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    prompt = f"""
Boss, I've reviewed the task: "{target_task['task']}"
Status: {'Completed' if target_task.get('completed') else 'Pending'}
Deadline: {target_task.get('deadline', 'None Set')}

COO, give me the master plan to execute this.
"""
    
    try:
        async with AI_SEMAPHORE:
            roadmap = await asyncio.get_event_loop().run_in_executor(
                None, generate_ai_response, prompt, COO_SYSTEM_PROMPT, 0.2, 1000, "COO"
            )
        return {"roadmap": roadmap or "Boss, I could not get a reliable COO briefing this time. Please try again in a moment."}
    except Exception as e:
        logger.error(f"COO Consultation Error: {e}")
        return {"error": "COO is in a meeting. Please try again later."}

@router.get("/stats")
def get_stats(request: Request):
    try:
        user_email = _get_user_email(request)
        tasks = get_tasks(user_email) if user_email else []
        return {
            "pending_tasks": len([t for t in tasks if not t['completed']]),
            "completed_tasks": len([t for t in tasks if t['completed']])
        }
    except: return {"pending_tasks": 0, "completed_tasks": 0}

@router.get("/tasks")
def fetch_tasks(request: Request):
    user_email = _get_user_email(request)
    if not user_email:
        logger.warning("Fetch tasks called without authenticated user email.")
        return {"tasks": []}
    
    tasks = get_tasks(user_email)
    logger.info(f"Fetch tasks API returning {len(tasks)} tasks for {user_email}")
    return {"tasks": tasks}

@router.get("/settings")
def get_settings_route(request: Request):
    user_email = _get_user_email(request)
    return {
        "fetch_limit": get_setting(user_email, "fetch_limit", "10"),
        "default_action": get_setting(user_email, "default_action", "draft"),
        "fetch_priority": get_setting(user_email, "fetch_priority", "all"),
        "ai_detail_level": get_setting(user_email, "ai_detail_level", "medium"),
        "ai_tone": get_setting(user_email, "ai_tone", "professional"),
        "ai_persona": get_setting(user_email, "ai_persona", ""),
        "ai_details": get_setting(user_email, "ai_details", ""),
        "university_mode": get_setting(user_email, "university_mode", "false") == "true",
        "university_roll_number": get_setting(user_email, "university_roll_number", ""),
        "university_semester": get_setting(user_email, "university_semester", ""),
        "university_section": get_setting(user_email, "university_section", ""),
        "university_course": get_setting(user_email, "university_course", ""),
        "university_specialization": get_setting(user_email, "university_specialization", ""),
        "university_campus": get_setting(user_email, "university_campus", ""),
        "university_ignore_other_sections": get_setting(user_email, "university_ignore_other_sections", "false") == "true",
        "university_filter_mode": get_setting(user_email, "university_filter_mode", "balanced"),
        "accent_color": get_setting(user_email, "accent_color", "#8b5cf6")
    }

@router.post("/settings")
def save_settings_route(settings: SettingsRequest, request: Request):
    user_email = _get_user_email(request)
    if not user_email:
        logger.error("Save settings failed: No authenticated user email found in request.")
        raise HTTPException(status_code=401, detail="Unauthorized")
    
    logger.info(f"Saving settings for {user_email}: Tone={settings.ai_tone}, PersonaLength={len(settings.ai_persona)}")
    
    # PERFORMANCE: Use a single DB session to persist all settings at once
    # This is much faster and more reliable than multiple individual calls
    from app.models import SessionLocal, Setting, db_write_lock
    
    with db_write_lock:
        db = SessionLocal()
        try:
            # Map of settings to persist
            payload = {
                "fetch_limit": str(settings.fetch_limit),
                "default_action": settings.default_action,
                "fetch_priority": settings.fetch_priority,
                "ai_detail_level": settings.ai_detail_level,
                "ai_tone": settings.ai_tone,
                "ai_persona": settings.ai_persona,
                "ai_details": settings.ai_details,
                "university_mode": str(settings.university_mode).lower(),
                "university_roll_number": settings.university_roll_number,
                "university_semester": settings.university_semester,
                "university_section": settings.university_section,
                "university_course": settings.university_course,
                "university_specialization": settings.university_specialization,
                "university_campus": settings.university_campus,
                "university_ignore_other_sections": str(settings.university_ignore_other_sections).lower(),
                "university_filter_mode": settings.university_filter_mode,
                "accent_color": settings.accent_color
            }

            for key, value in payload.items():
                existing = db.query(Setting).filter(Setting.user_email == user_email, Setting.key == key).first()
                if existing:
                    existing.value = str(value)
                else:
                    db.add(Setting(user_email=user_email, key=key, value=str(value)))
            
            db.commit()
            # Invalidate the lru_cache for get_setting
            from app.services.task_store import get_setting
            get_setting.cache_clear()
            
            logger.info(f"Settings successfully persisted for {user_email}")
            return {"message": "Settings saved"}
        except Exception as e:
            db.rollback()
            logger.error(f"Failed to persist settings for {user_email}: {e}")
            raise HTTPException(status_code=500, detail=f"Persistence Error: {str(e)}")
        finally:
            db.close()

@router.post("/logout")
def logout(response: Response):
    from app.services.gmail_service import TOKEN_PATH
    if os.path.exists(TOKEN_PATH):
        try: os.remove(TOKEN_PATH)
        except: pass
    response.delete_cookie("mailpilot_token", path="/", secure=_cookie_secure(), samesite="lax")
    return {"message": "OK"}

@router.get("/oauth2callback")
def oauth2callback_redirect(code: str, state: str = None):
    """Google redirects here via GET. Forward to frontend SPA."""
    from urllib.parse import urlencode
    from starlette.responses import RedirectResponse
    return RedirectResponse(url=f"/?{urlencode({'code': code, 'state': state or ''})}", status_code=302)

@router.post("/oauth2callback")
def oauth2callback(req: AuthCallbackRequest, response: Response):
    try:
        creds = exchange_code(req.code, req.state, req.code_verifier)
        # Fetch email once during login to cache it in the cookie
        from googleapiclient.discovery import build
        service = build('gmail', 'v1', credentials=creds)
        info = get_user_info(service)
        email = info.get('email') if info else None
        
        _set_token_cookie(response, creds, email)
        return {"message": "OK"}
    except Exception as e:
        logger.error(f"OAuth Callback Error: {e}")
        raise HTTPException(status_code=400, detail=f"OAuth Error: {str(e)}")

@router.put("/complete-task/{task_id}")
def complete_task(task_id: Union[int, str], request: Request):
    try: tid = int(task_id)
    except ValueError: raise HTTPException(status_code=400, detail="Invalid task ID")
    user_email = _get_user_email(request)
    if mark_task_complete(tid, user_email): return {"message": "Success"}
    raise HTTPException(status_code=404, detail="Task not found")

@router.delete("/delete-task/{task_id}")
def remove_task(task_id: Union[int, str], request: Request):
    try: tid = int(task_id)
    except ValueError: raise HTTPException(status_code=400, detail="Invalid task ID")
    user_email = _get_user_email(request)
    if delete_task(tid, user_email): return {"message": "Success"}
    raise HTTPException(status_code=404, detail="Task not found")

@router.delete("/tasks/clear")
def clear_tasks(request: Request):
    user_email = _get_user_email(request)
    if not user_email:
        raise HTTPException(status_code=401, detail="Unauthorized")
    clear_all_tasks(user_email)
    return {"message": "Success"}

@router.post("/generate-reply")
async def generate_reply_route(req: ReplyRequest, request: Request):
    user_email = _get_user_email(request)
    from app.services.ai_service import generate_reply
    from app.services.task_store import get_setting
    
    # Personalize the draft using user's saved preferences
    # ALLOW OVERRIDE: If the request specifies a tone (e.g. from the UI refinement buttons), use it.
    tone = req.tone if hasattr(req, 'tone') and req.tone else get_setting(user_email, "ai_tone", "professional")
    details = req.ai_details if req.ai_details is not None else get_setting(user_email, "ai_details", "") 
    
    # FETCH STYLE MEMORY
    from app.services.task_store import get_style_examples
    style_examples = get_style_examples(user_email, 5)
    
    reply = await asyncio.to_thread(
        generate_reply,
        req.content,
        tone=tone,
        user_identifier=user_email,
        user_details=details,
        style_examples=style_examples
    )
    return {"reply": reply}

@router.post("/send-email")
async def send_email_route(req: SendEmailRequest, request: Request):
    try:
        send_email_direct(req.to, req.subject, req.body, _get_cookie(request))
        
        # SAVE TO STYLE MEMORY (For future drafting learning)
        user_email = _get_user_email(request)
        if user_email:
            from app.services.task_store import add_style_reference
            add_style_reference(user_email, req.body)
            
        return {"message": "Email sent"}
    except Exception as e:
        logger.error(f"Send Email Error: {e}")
        raise HTTPException(status_code=500, detail="Failed to send email")

@router.post("/extract-tasks")
async def extract_tasks_route(req: dict, request: Request):
    async with AI_SEMAPHORE:
        user_email = _get_user_email(request)
        content = req.get("email", "")
        
        # Load feedback even for manual extraction to improve quality
        task_pos, task_neg = ([], [])
        if user_email:
            from app.services.task_store import get_task_feedback_examples
            task_pos, task_neg = get_task_feedback_examples(user_email, 5)
            
        persona = _build_effective_persona(user_email) if user_email else ""
        tasks = extract_tasks(content, pos_examples=task_pos, neg_examples=task_neg, user_persona=persona)
        
        if user_email: 
            # Manual tasks are also saved to the dashboard for persistence
            final_tasks = save_tasks(tasks, user_email, message_id="manual-compose", source_snippet=content)
            return {"tasks": final_tasks}
        return {"tasks": tasks}

@router.post("/archive-thread/{thread_id}")
async def archive_thread_route(thread_id: str, request: Request):
    from app.services.gmail_service import archive_thread
    success = archive_thread(thread_id, _get_cookie(request))
    if not success:
        raise HTTPException(status_code=500, detail="Failed to archive thread")
    return {"message": "Thread archived"}
